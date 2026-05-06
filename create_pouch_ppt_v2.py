from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Palette ────────────────────────────────────────────────────
C_BG      = RGBColor(0x0B, 0x14, 0x26)
C_CARD    = RGBColor(0x16, 0x24, 0x3A)
C_CARD2   = RGBColor(0x1E, 0x2F, 0x4A)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_MUTED   = RGBColor(0x94, 0xA3, 0xB8)
C_LIGHT   = RGBColor(0xCB, 0xD5, 0xE1)

C_BLUE    = RGBColor(0x38, 0xBD, 0xF8)   # sky  — 3 Side Seal
C_GREEN   = RGBColor(0x34, 0xD3, 0x99)   # emerald — Center Seal
C_ORANGE  = RGBColor(0xFB, 0x92, 0x3C)   # orange — Standup/Zipper
C_PURPLE  = RGBColor(0xA7, 0x8B, 0xFA)   # violet — Both Gusset
C_TEAL    = RGBColor(0x2D, 0xD4, 0xBF)   # teal — 3D Flat
C_YELLOW  = RGBColor(0xFB, 0xBF, 0x24)   # amber
C_RED     = RGBColor(0xF8, 0x71, 0x71)   # red

POUCH_C = [C_BLUE, C_GREEN, C_ORANGE, C_PURPLE, C_TEAL]

# ── Base helpers ───────────────────────────────────────────────
def bg(sl, col=C_BG):
    s = sl.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    s.fill.solid(); s.fill.fore_color.rgb = col; s.line.fill.background()

def rect(sl, x, y, w, h, fill, lc=None, lw=Pt(0)):
    s = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if lc: s.line.color.rgb = lc; s.line.width = lw
    else: s.line.fill.background()
    return s

def txt(sl, t, x, y, w, h, sz=14, bold=False, col=C_WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = t
    r.font.size = Pt(sz); r.font.bold = bold
    r.font.color.rgb = col; r.font.italic = italic
    return tb

def top_bar(sl, col):
    rect(sl, 0, 0, 13.33, 0.07, col)

def bot_bar(sl, col=C_BLUE):
    rect(sl, 0, 7.38, 13.33, 0.12, col)

def slide_num(sl, n, total=14):
    txt(sl, f"{n} / {total}", 12.2, 7.18, 1.0, 0.28, sz=9,
        col=C_MUTED, align=PP_ALIGN.RIGHT)

def section_tag(sl, label, col):
    rect(sl, 0.5, 0.18, len(label)*0.13+0.4, 0.38, col)
    txt(sl, label, 0.52, 0.18, len(label)*0.13+0.36, 0.38,
        sz=10, bold=True, col=C_BG, align=PP_ALIGN.CENTER)

def heading(sl, h1, h2="", col=C_WHITE):
    txt(sl, h1, 0.5, 0.65, 12.5, 0.55, sz=30, bold=True, col=col)
    if h2:
        txt(sl, h2, 0.5, 1.18, 12.5, 0.4, sz=15, col=C_MUTED)
    rect(sl, 0.5, 1.62, 12.33, 0.04, col)

def pill(sl, label, val, x, y, col, w=2.55, h=0.72):
    rect(sl, x, y, w, h, C_CARD2)
    rect(sl, x, y, 0.06, h, col)
    txt(sl, label, x+0.14, y+0.04, w-0.2, 0.26, sz=8, bold=True, col=col)
    txt(sl, val,   x+0.14, y+0.3,  w-0.2, 0.36, sz=12, bold=True, col=C_WHITE)

# ══════════════════════════════════════════════════════════════
# SLIDE 1  ── TITLE
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
# gradient strip
rect(sl, 0, 3.0, 13.33, 4.5, RGBColor(0x08, 0x10, 0x1E))
rect(sl, 0, 0, 13.33, 0.07, C_BLUE)
# logo badge
rect(sl, 0.5, 0.22, 2.6, 0.5, C_BLUE)
txt(sl, "AJ SHRINK ERP", 0.5, 0.22, 2.6, 0.5, sz=12, bold=True,
    col=C_BG, align=PP_ALIGN.CENTER)
# main title
txt(sl, "Flexible Packaging", 1.8, 1.25, 9.8, 0.75,
    sz=36, col=C_MUTED, align=PP_ALIGN.CENTER)
txt(sl, "Pouch Types & Their Applications", 0.9, 1.95, 11.5, 1.0,
    sz=50, bold=True, col=C_WHITE, align=PP_ALIGN.CENTER)
txt(sl, "Where they are used  ·  What goes inside  ·  Film & Cylinder Calculations",
    1.5, 3.05, 10.3, 0.5, sz=16, col=C_BLUE, align=PP_ALIGN.CENTER, italic=True)
rect(sl, 4.8, 3.7, 3.7, 0.05, C_BLUE)
# 5 type pills
names = ["3 Side Seal","Center Seal","Standup / Zipper","Both Gusset","3D Flat Bottom"]
for i,(n,c) in enumerate(zip(names, POUCH_C)):
    xi = 0.65 + i*2.41
    rect(sl, xi, 4.1, 2.22, 0.52, c)
    txt(sl, n, xi, 4.1, 2.22, 0.52, sz=10, bold=True,
        col=C_BG, align=PP_ALIGN.CENTER)
txt(sl, "Client Presentation  —  Confidential",
    4.0, 6.95, 5.5, 0.3, sz=10, col=C_MUTED,
    align=PP_ALIGN.CENTER, italic=True)
bot_bar(sl, C_BLUE)

# ══════════════════════════════════════════════════════════════
# SLIDE 2  ── WHY POUCH TYPE MATTERS
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl); top_bar(sl, C_YELLOW)
section_tag(sl, "INTRODUCTION", C_YELLOW)
heading(sl, "Why Choosing the Right Pouch Matters",
        "The pouch type directly controls print area, material cost, and cylinder size", C_YELLOW)

points = [
    (C_BLUE,   "Film Width",
     "Every pouch type adds SEALS or GUSSETS to the base width.\n"
     "This determines how many pouches fit ACROSS the film roll (AC UPS).\n"
     "More UPS across = lower cost per unit."),
    (C_GREEN,  "Cylinder Circumference",
     "Seals + gussets added to height decide the exact cylinder circumference.\n"
     "Wrong measurement = cylinder won't divide evenly = print misalignment.\n"
     "Every mm counts for a ₹50,000+ cylinder order."),
    (C_ORANGE, "Material Waste",
     "Incorrect seal allowance wastes film with every meter printed.\n"
     "1mm error × 5,00,000 pouches = 500 meters of wasted film.\n"
     "Proper planning eliminates dead margin and side waste."),
    (C_PURPLE, "Structure for the Product",
     "Heavier products need gussets. Liquids need tight center seals.\n"
     "Standing products need standup gussets. Premium products need 3D shape.\n"
     "Wrong structure = pouch fails on the filling line."),
]
for i,(col,title,desc) in enumerate(points):
    px = 0.5 + (i%2)*6.42
    py = 1.8 + (i//2)*2.55
    rect(sl, px, py, 6.2, 2.38, C_CARD)
    rect(sl, px, py, 0.07, 2.38, col)
    rect(sl, px+0.18, py+0.06, 0.45, 0.45, col)
    txt(sl, str(i+1), px+0.18, py+0.06, 0.45, 0.45,
        sz=18, bold=True, col=C_BG, align=PP_ALIGN.CENTER)
    txt(sl, title, px+0.74, py+0.1, 5.3, 0.42, sz=14, bold=True, col=col)
    txt(sl, desc,  px+0.18, py+0.6, 5.9, 1.65, sz=10, col=C_LIGHT, wrap=True)

slide_num(sl, 2); bot_bar(sl, C_YELLOW)

# ══════════════════════════════════════════════════════════════
# SLIDE 3  ── QUICK REFERENCE TABLE
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl); top_bar(sl, C_BLUE)
section_tag(sl, "OVERVIEW", C_BLUE)
heading(sl, "All 5 Pouch Types at a Glance",
        "Film Width formula  ·  Repeat Length formula  ·  Typical products")

rows = [
    (C_BLUE,   "①", "3 Side Seal",
     "W + 2×Side Seal",
     "H + Top + Bottom Seal",
     "Spices, sachets, namkeen, seeds, sugar, flour"),
    (C_GREEN,  "②", "Center Seal",
     "2W + Center Seal",
     "H + Top + Bottom Seal",
     "Milk pouches, oil pouches, liquid sachets, shampoo"),
    (C_ORANGE, "③", "Standup / Zipper",
     "W + 2×Side Seal",
     "H + Top + Gusset÷2",
     "Dry fruits, pet food, coffee, protein powder, snacks"),
    (C_PURPLE, "④", "Both Side Gusset",
     "W + 2×Side Gusset",
     "H + Top + Bottom Seal",
     "Rice, atta, sugar (5kg+), cement additives, animal feed"),
    (C_TEAL,   "⑤", "3D / Flat Bottom",
     "W + 2×Side Gusset",
     "H + Top + Btm Gusset÷2",
     "Premium coffee, tea, cosmetics, organic products"),
]
hdr_y = 1.75
rect(sl, 0.5, hdr_y, 12.33, 0.42, RGBColor(0x10, 0x1C, 0x30))
for hdr,cx,cw in zip(
    ["#","Pouch Type","Film Width","Repeat Length","Typical Products"],
    [0.5, 1.1, 3.3, 6.1, 8.9],
    [0.5, 2.1, 2.7, 2.7, 4.0]):
    txt(sl, hdr, cx+0.06, hdr_y+0.06, cw, 0.3,
        sz=9, bold=True, col=C_BLUE, align=PP_ALIGN.LEFT)

for ri,(col,num,name,fw,rep,prod) in enumerate(rows):
    ry = hdr_y+0.48 + ri*1.02
    fill = RGBColor(0x10,0x1C,0x30) if ri%2==0 else C_CARD
    rect(sl, 0.5, ry, 12.33, 0.95, fill)
    rect(sl, 0.5, ry, 0.07, 0.95, col)
    txt(sl, num,  0.52, ry+0.2, 0.5, 0.55, sz=18, bold=True, col=col, align=PP_ALIGN.CENTER)
    txt(sl, name, 1.1,  ry+0.06, 2.0, 0.82, sz=12, bold=True, col=C_WHITE)
    txt(sl, fw,   3.3,  ry+0.06, 2.6, 0.82, sz=10, col=C_LIGHT)
    txt(sl, rep,  6.1,  ry+0.06, 2.6, 0.82, sz=10, col=C_LIGHT)
    txt(sl, prod, 8.9,  ry+0.06, 3.85,0.82, sz=9,  col=C_MUTED, italic=True)

slide_num(sl, 3); bot_bar(sl, C_BLUE)

# ══════════════════════════════════════════════════════════════
# Helper: full detail slide per pouch type
# ══════════════════════════════════════════════════════════════
def pouch_slide(prs, slide_n, num, name, col,
                tagline,
                film_f, repeat_f, example_calc,
                uses,         # list of strings
                products,     # list of strings
                materials,    # list of strings
                diagram_fn,
                pro_tip=""):

    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl); top_bar(sl, col)

    # Number badge + title
    rect(sl, 0.5, 0.14, 0.48, 0.48, col)
    txt(sl, num, 0.5, 0.14, 0.48, 0.48, sz=20, bold=True,
        col=C_BG, align=PP_ALIGN.CENTER)
    txt(sl, name.upper(), 1.08, 0.12, 11.0, 0.52, sz=26, bold=True, col=col)
    txt(sl, tagline, 1.08, 0.62, 11.0, 0.35, sz=12, col=C_MUTED, italic=True)
    rect(sl, 0.5, 1.02, 12.33, 0.04, col)

    # ── LEFT column (diagram + formula + example) ─────────────
    lw = 5.4  # left col width

    # Diagram area
    rect(sl, 0.5, 1.12, lw, 3.1, C_CARD)
    diagram_fn(sl, col, 0.5, 1.12, lw, 3.1)

    # Formula box
    rect(sl, 0.5, 4.28, lw, 1.38, C_CARD2)
    rect(sl, 0.5, 4.28, 0.07, 1.38, col)
    txt(sl, "FILM WIDTH (across web):", 0.65, 4.32, lw-0.2, 0.28,
        sz=8, bold=True, col=col)
    txt(sl, film_f, 0.65, 4.58, lw-0.2, 0.38, sz=13, bold=True, col=C_WHITE)
    rect(sl, 0.65, 4.98, lw-0.22, 0.03, RGBColor(0x25,0x38,0x55))
    txt(sl, "REPEAT LENGTH (cyl. circumference):", 0.65, 5.02, lw-0.2, 0.28,
        sz=8, bold=True, col=col)
    txt(sl, repeat_f, 0.65, 5.28, lw-0.2, 0.35, sz=13, bold=True, col=C_WHITE)

    # Example calc chip
    rect(sl, 0.5, 5.7, lw, 0.55, RGBColor(0x0A,0x1F,0x14))
    rect(sl, 0.5, 5.7, 0.07, 0.55, C_GREEN)
    txt(sl, "Example: " + example_calc, 0.65, 5.74, lw-0.2, 0.44,
        sz=10, col=C_GREEN, bold=True)

    # ── RIGHT column (uses / products / materials) ─────────────
    rx = 6.15; rw = 6.68

    # WHERE IS IT USED
    rect(sl, rx, 1.12, rw, 1.85, C_CARD)
    rect(sl, rx, 1.12, 0.07, 1.85, col)
    txt(sl, "WHERE IS IT USED?", rx+0.18, 1.16, rw-0.28, 0.3,
        sz=9, bold=True, col=col)
    for ui,u in enumerate(uses):
        txt(sl, "▸  "+u, rx+0.18, 1.48+ui*0.32, rw-0.3, 0.3,
            sz=10, col=C_LIGHT)

    # WHAT PRODUCTS GO IN
    rect(sl, rx, 3.04, rw, 1.68, C_CARD2)
    rect(sl, rx, 3.04, 0.07, 1.68, C_YELLOW)
    txt(sl, "WHAT PRODUCTS GO INSIDE?", rx+0.18, 3.08, rw-0.28, 0.3,
        sz=9, bold=True, col=C_YELLOW)
    # 2-col product list
    half = len(products)//2 + len(products)%2
    for pi,p in enumerate(products):
        col_off = 0 if pi < half else (rw-0.3)/2
        row_off = pi if pi < half else pi-half
        txt(sl, "• "+p, rx+0.18+col_off, 3.4+row_off*0.3, (rw-0.3)/2, 0.28,
            sz=9, col=C_MUTED)

    # MATERIAL / FILM STRUCTURE
    rect(sl, rx, 4.78, rw, 1.5, C_CARD)
    rect(sl, rx, 4.78, 0.07, 1.5, C_PURPLE)
    txt(sl, "TYPICAL FILM STRUCTURES USED", rx+0.18, 4.82, rw-0.28, 0.3,
        sz=9, bold=True, col=C_PURPLE)
    for mi,m in enumerate(materials):
        txt(sl, f"  {mi+1}.  {m}", rx+0.18, 5.14+mi*0.3, rw-0.3, 0.28,
            sz=9, col=C_LIGHT)

    # Pro tip
    if pro_tip:
        rect(sl, rx, 6.34, rw, 0.62, RGBColor(0x0C,0x1F,0x14))
        rect(sl, rx, 6.34, 0.07, 0.62, C_GREEN)
        txt(sl, "💡  " + pro_tip, rx+0.18, 6.38, rw-0.28, 0.52,
            sz=9, col=C_GREEN, italic=True)

    slide_num(sl, slide_n); bot_bar(sl, col)
    return sl

# ══════════════════════════════════════════════════════════════
# Diagram functions
# ══════════════════════════════════════════════════════════════
def d_threeside(sl, col, bx, by, bw, bh):
    mx = bx+bw/2; my = by+bh*0.48
    pw=1.7; ph=2.1; ss=0.25; ts=0.2; bs=0.18
    ox=mx-(pw/2+ss); oy=my-ph/2-ts/2
    fw=pw+2*ss; fh=ph+ts+bs
    # full film bg
    rect(sl, ox, oy, fw, fh, RGBColor(0x0C,0x22,0x16))
    # seals
    rect(sl, ox, oy, ss, fh, RGBColor(0x14,0x55,0x30))
    rect(sl, ox+ss+pw, oy, ss, fh, RGBColor(0x14,0x55,0x30))
    rect(sl, ox+ss, oy, pw, ts, RGBColor(0x14,0x55,0x30))
    rect(sl, ox+ss, oy+ts+ph, pw, bs, RGBColor(0x14,0x55,0x30))
    rect(sl, ox+ss, oy+ts, pw, ph, RGBColor(0x17,0x3A,0x26))
    # labels
    txt(sl,"BODY",     mx-0.3, my-0.12, 0.6, 0.25, sz=8, bold=True, col=C_WHITE, align=PP_ALIGN.CENTER)
    txt(sl,"SIDE SEAL",ox-0.04,oy+fh*0.38,ss+0.08,0.28,sz=5,col=col,align=PP_ALIGN.CENTER)
    txt(sl,"TOP SEAL", ox+ss, oy-0.02, pw, ts+0.04, sz=6, col=col, align=PP_ALIGN.CENTER)
    txt(sl,"BTM SEAL", ox+ss, oy+ts+ph-0.01, pw, bs+0.06, sz=6, col=col, align=PP_ALIGN.CENTER)
    # fold line
    rect(sl, ox+ss, oy, 0.03, fh, RGBColor(0x0E,0x6B,0x40))
    rect(sl, ox+ss+pw-0.03, oy, 0.03, fh, RGBColor(0x0E,0x6B,0x40))
    # dimension arrows
    txt(sl, f"← W+2S = Film Width →", bx+0.1, oy+fh+0.08, bw-0.2, 0.25,
        sz=8, bold=True, col=col, align=PP_ALIGN.CENTER)
    txt(sl, "3 SIDES SEALED: Left + Right + Bottom", bx+0.1, oy+fh+0.32, bw-0.2, 0.25,
        sz=8, col=C_MUTED, align=PP_ALIGN.CENTER, italic=True)

def d_centerseal(sl, col, bx, by, bw, bh):
    mx=bx+bw/2; my=by+bh*0.48
    hw=0.9; cs=0.2; ph=2.1; ts=0.2; bs=0.18
    ox=mx-hw-cs/2; oy=my-ph/2-ts/2
    fh=ph+ts+bs
    rect(sl,ox,oy,hw,fh,RGBColor(0x18,0x16,0x3E))
    rect(sl,ox+hw+cs,oy,hw,fh,RGBColor(0x18,0x16,0x3E))
    rect(sl,ox+hw,oy,cs,fh,RGBColor(0x44,0x1A,0x80))
    rect(sl,ox,oy,hw*2+cs,ts,RGBColor(0x35,0x1B,0x76))
    rect(sl,ox,oy+ts+ph,hw*2+cs,bs,RGBColor(0x35,0x1B,0x76))
    txt(sl,"A",ox+hw*0.38,my-0.12,0.3,0.28,sz=10,bold=True,col=C_WHITE,align=PP_ALIGN.CENTER)
    txt(sl,"B",ox+hw+cs+hw*0.32,my-0.12,0.3,0.28,sz=10,bold=True,col=C_WHITE,align=PP_ALIGN.CENTER)
    txt(sl,"CENTER\nSEAL",ox+hw,oy+fh*0.35,cs+0.04,0.45,sz=5,col=col,align=PP_ALIGN.CENTER)
    txt(sl,"TOP SEAL",ox,oy-0.02,hw*2+cs,ts+0.04,sz=6,col=col,align=PP_ALIGN.CENTER)
    txt(sl,"BTM SEAL",ox,oy+ts+ph-0.01,hw*2+cs,bs+0.05,sz=6,col=col,align=PP_ALIGN.CENTER)
    txt(sl,"← 2W + Center Seal = Film Width →",bx+0.1,oy+fh+0.08,bw-0.2,0.25,
        sz=8,bold=True,col=col,align=PP_ALIGN.CENTER)
    txt(sl,"Film wraps around — A & B are same printed panel",bx+0.1,oy+fh+0.32,bw-0.2,0.25,
        sz=8,col=C_MUTED,align=PP_ALIGN.CENTER,italic=True)

def d_standup(sl, col, bx, by, bw, bh):
    mx=bx+bw/2; my=by+bh*0.43
    pw=1.7; ph=1.9; ss=0.24; ts=0.18; gus=0.4
    ox=mx-(pw/2+ss); oy=my-ph/2-ts/2
    fw=pw+2*ss; fh=ph+ts+gus
    rect(sl,ox,oy,fw,fh,RGBColor(0x18,0x1A,0x38))
    rect(sl,ox,oy,ss,fh-gus,RGBColor(0x2A,0x27,0x62))
    rect(sl,ox+ss+pw,oy,ss,fh-gus,RGBColor(0x2A,0x27,0x62))
    rect(sl,ox+ss,oy,pw,ts,RGBColor(0x46,0x42,0x96))
    rect(sl,ox+ss,oy+ts,pw,ph,RGBColor(0x20,0x24,0x4E))
    rect(sl,ox,oy+ts+ph,fw,gus,RGBColor(0x08,0x72,0x88))
    txt(sl,"BODY",mx-0.3,my+0.0,0.6,0.25,sz=8,bold=True,col=C_WHITE,align=PP_ALIGN.CENTER)
    txt(sl,"GUSSET",mx-0.38,oy+ts+ph+0.06,0.85,0.22,sz=7,bold=True,col=C_BG,align=PP_ALIGN.CENTER)
    txt(sl,"TOP SEAL",ox+ss,oy-0.02,pw,ts+0.04,sz=6,col=col,align=PP_ALIGN.CENTER)
    txt(sl,"SIDE",ox-0.02,oy+0.45,ss+0.05,0.22,sz=5,col=col,align=PP_ALIGN.CENTER)
    rect(sl,ox+ss/2-0.015,oy,0.03,fh-gus,RGBColor(0x60,0x56,0xD0))
    rect(sl,ox+ss+pw+ss/2-0.015,oy,0.03,fh-gus,RGBColor(0x60,0x56,0xD0))
    txt(sl,"← W+2×Side Seal = Film Width →",bx+0.1,oy+fh+0.08,bw-0.2,0.25,
        sz=8,bold=True,col=col,align=PP_ALIGN.CENTER)
    txt(sl,"Gusset÷2 folds into repeat length",bx+0.1,oy+fh+0.32,bw-0.2,0.25,
        sz=8,col=C_MUTED,align=PP_ALIGN.CENTER,italic=True)

def d_bothgusset(sl, col, bx, by, bw, bh):
    mx=bx+bw/2; my=by+bh*0.46
    pw=1.5; ph=2.1; sg=0.38; ts=0.18; bs=0.18
    ox=mx-(pw/2+sg); oy=my-ph/2-ts/2
    fw=pw+2*sg; fh=ph+ts+bs
    rect(sl,ox,oy,fw,fh,RGBColor(0x1C,0x14,0x28))
    rect(sl,ox,oy,sg,fh,RGBColor(0x52,0x1B,0x80))
    rect(sl,ox+sg+pw,oy,sg,fh,RGBColor(0x52,0x1B,0x80))
    rect(sl,ox+sg,oy,pw,ts,RGBColor(0x70,0x35,0xD0))
    rect(sl,ox+sg,oy+ts+ph,pw,bs,RGBColor(0x70,0x35,0xD0))
    rect(sl,ox+sg,oy+ts,pw,ph,RGBColor(0x28,0x18,0x3C))
    # fold lines
    rect(sl,ox+sg/2-0.015,oy,0.03,fh,RGBColor(0x98,0x7A,0xF2))
    rect(sl,ox+sg+pw+sg/2-0.015,oy,0.03,fh,RGBColor(0x98,0x7A,0xF2))
    txt(sl,"BODY",mx-0.28,my+0.05,0.6,0.25,sz=8,bold=True,col=C_WHITE,align=PP_ALIGN.CENTER)
    txt(sl,"SIDE\nGUSSET",ox+0.02,oy+fh*0.38,sg+0.02,0.38,sz=5,col=C_BG,align=PP_ALIGN.CENTER)
    txt(sl,"TOP SEAL",ox+sg,oy-0.02,pw,ts+0.04,sz=6,col=col,align=PP_ALIGN.CENTER)
    txt(sl,"BTM SEAL",ox+sg,oy+ts+ph-0.01,pw,bs+0.05,sz=6,col=col,align=PP_ALIGN.CENTER)
    txt(sl,"← W + 2×Side Gusset = Film Width →",bx+0.1,oy+fh+0.08,bw-0.2,0.25,
        sz=8,bold=True,col=col,align=PP_ALIGN.CENTER)
    txt(sl,"Gussets expand sideways for volume",bx+0.1,oy+fh+0.32,bw-0.2,0.25,
        sz=8,col=C_MUTED,align=PP_ALIGN.CENTER,italic=True)

def d_flatbottom(sl, col, bx, by, bw, bh):
    mx=bx+bw/2; my=by+bh*0.44
    pw=1.5; ph=1.9; sg=0.36; ts=0.18; bg2=0.44
    ox=mx-(pw/2+sg); oy=my-ph/2-ts/2
    fw=pw+2*sg; fh=ph+ts+bg2/2
    rect(sl,ox,oy,fw,fh,RGBColor(0x08,0x1C,0x14))
    rect(sl,ox,oy,sg,fh,RGBColor(0x04,0x48,0x2C))
    rect(sl,ox+sg+pw,oy,sg,fh,RGBColor(0x04,0x48,0x2C))
    rect(sl,ox+sg,oy,pw,ts,RGBColor(0x0A,0x70,0x46))
    rect(sl,ox+sg,oy+ts,pw,ph,RGBColor(0x0A,0x26,0x18))
    rect(sl,ox,oy+ts+ph,fw,bg2/2,RGBColor(0x0A,0x70,0x46))
    rect(sl,ox+sg,oy+ts+ph+bg2*0.23,pw,0.025,RGBColor(0x60,0xD8,0xA8))
    # fold lines
    rect(sl,ox+sg/2-0.015,oy,0.03,fh,RGBColor(0x28,0xC4,0x88))
    rect(sl,ox+sg+pw+sg/2-0.015,oy,0.03,fh,RGBColor(0x28,0xC4,0x88))
    txt(sl,"BODY",mx-0.28,my+0.02,0.6,0.25,sz=8,bold=True,col=C_WHITE,align=PP_ALIGN.CENTER)
    txt(sl,"SIDE\nGUSSET",ox+0.02,oy+fh*0.36,sg+0.02,0.36,sz=5,col=C_BG,align=PP_ALIGN.CENTER)
    txt(sl,"TOP SEAL",ox+sg,oy-0.02,pw,ts+0.04,sz=6,col=col,align=PP_ALIGN.CENTER)
    txt(sl,"BTM FOLD",ox+sg,oy+ts+ph+0.02,pw,0.22,sz=6,col=C_BG,align=PP_ALIGN.CENTER)
    txt(sl,"← W + 2×Side Gusset = Film Width →",bx+0.1,oy+fh+0.08,bw-0.2,0.25,
        sz=8,bold=True,col=col,align=PP_ALIGN.CENTER)
    txt(sl,"3D block shape — stands + has flat base",bx+0.1,oy+fh+0.32,bw-0.2,0.25,
        sz=8,col=C_MUTED,align=PP_ALIGN.CENTER,italic=True)

# ══════════════════════════════════════════════════════════════
# SLIDES 4–8: Individual Pouch Slides
# ══════════════════════════════════════════════════════════════

pouch_slide(prs, 4, "①", "3 Side Seal Pouch", C_BLUE,
    "Simplest pouch — three sides heat-sealed, one side fold",
    film_f   = "Film Width  =  W  +  (2 × Side Seal)",
    repeat_f = "Repeat  =  H  +  Top Seal  +  Bottom Seal",
    example_calc = "W=130, H=200, Top=10, Btm=8, Side=10 → Film=150mm, Repeat=218mm",
    uses=[
        "Small retail pouches — 10g to 500g range",
        "Single-serve sachets — sugar, salt, spice, sauce",
        "Pharma — tablet strips, powder sachets",
        "Agricultural — seed packets, fertilizer sachets",
        "Industrial — adhesive sachets, lubricant pouches",
    ],
    products=[
        "Namkeen / chips", "Spices (haldi, mirchi)", "Salt & sugar sachets",
        "Tea & coffee sachets", "Cement additives", "Seeds",
        "Sauce / ketchup", "Tablet pouches", "Detergent powder", "Seeds",
    ],
    materials=[
        "BOPP + CPP  (most common for dry products)",
        "PET + PE  (higher barrier, good for moisture)",
        "PET + AL Foil + PE  (maximum barrier — pharma, spices)",
        "BOPP + MET BOPP  (shiny look — premium snacks)",
    ],
    diagram_fn=d_threeside,
    pro_tip="Most economical pouch — highest AC UPS possible. Side seal width 8–12mm typical. Use for anything dry and non-premium."
)

pouch_slide(prs, 5, "②", "Center Seal Pouch", C_GREEN,
    "Film wraps around the product — sealed at the back center",
    film_f   = "Film Width  =  (W × 2)  +  Center Seal Width",
    repeat_f = "Repeat  =  H  +  Top Seal  +  Bottom Seal",
    example_calc = "W=80 (half), H=200, Top=10, Btm=8, Seal=12 → Film=172mm, Repeat=218mm",
    uses=[
        "Liquid pouches — milk, oil, drinking water",
        "Flow-wrap / pillow pack on HFFS machines",
        "Loose items — candy, biscuits, small hardware",
        "Bulk packaging — continuous film roll",
        "Instant noodles, bread — bakery flow wrap",
    ],
    products=[
        "Milk (200ml, 500ml)", "Refined oil", "Drinking water",
        "Candy & toffee", "Biscuits", "Bread loaf",
        "Shampoo sachet", "Tomato sauce", "Noodles", "Ice cream",
    ],
    materials=[
        "PET + PE  (standard liquid barrier — milk, oil)",
        "LLDPE single layer  (low-cost water pouches)",
        "PET + MET PET + PE  (oil, shampoo — high barrier)",
        "BOPP + CPP  (flow wrap — biscuits, candy)",
    ],
    diagram_fn=d_centerseal,
    pro_tip="W entered is HALF the pouch face. Film Width = 2W+seal. HFFS machines use this structure — continuous film fed from roll."
)

pouch_slide(prs, 6, "③", "Standup Pouch / Zipper Pouch", C_ORANGE,
    "Bottom gusset lets the pouch stand upright on shelf",
    film_f   = "Film Width  =  W  +  (2 × Side Seal)",
    repeat_f = "Repeat  =  H  +  Top Seal  +  Bottom Gusset ÷ 2",
    example_calc = "W=150, H=220, Top=12, Side=10, Gusset=60 → Film=170mm, Repeat=262mm",
    uses=[
        "Retail shelf display — stands without support",
        "Premium FMCG — dry fruits, nuts, protein powder",
        "Pet food — kibble, treats, supplements",
        "E-commerce packaging — reusable with zipper",
        "Organic / health food — premium segment",
    ],
    products=[
        "Dry fruits & nuts", "Protein powder", "Pet food & treats",
        "Coffee beans", "Quinoa & superfoods", "Washing powder",
        "Gummies & supplements", "Baby food (puffed)", "Makhana", "Seeds",
    ],
    materials=[
        "PET + AL Foil + PE  (best barrier — coffee, protein)",
        "BOPP + CPP  (cost effective — snacks, dry items)",
        "PET + MET PET + PE  (metallic look — premium)",
        "PET + PE with matte OPP  (organic / natural look)",
    ],
    diagram_fn=d_standup,
    pro_tip="Gusset/2 added to repeat — e.g. 60mm gusset adds 30mm to cylinder. Zipper pouch same as standup + zipper strip cut 20–25mm from top."
)

pouch_slide(prs, 7, "④", "Both Side Gusset Pouch", C_PURPLE,
    "Gussets on both sides expand laterally for high-volume products",
    film_f   = "Film Width  =  W  +  (2 × Side Gusset)",
    repeat_f = "Repeat  =  H  +  Top Seal  +  Bottom Seal",
    example_calc = "W=200, H=300, Top=15, Btm=12, Gusset=40 → Film=280mm, Repeat=327mm",
    uses=[
        "Heavy grain packaging — 2kg, 5kg, 10kg bags",
        "Industrial products — chemicals, fertilizer, animal feed",
        "Bulk retail — atta, rice, dal, sugar",
        "Construction — tile adhesive, grout, putty",
        "Agriculture — seeds, pesticide powder",
    ],
    products=[
        "Rice (2kg, 5kg, 10kg)", "Wheat flour (atta)", "Sugar & jaggery",
        "Dal & pulses", "Poha & dry grains", "Animal feed",
        "Fertilizer (IFFCO)", "Tile adhesive powder", "Cement admixture", "Seeds (bulk)",
    ],
    materials=[
        "BOPP + CPP  (standard grain bags — rice, atta)",
        "Woven PP laminate  (heavy-duty — fertilizer, cement)",
        "PET + AL + PE  (moisture-sensitive — chemicals)",
        "MET BOPP + CPP  (shiny grain bags — premium rice)",
    ],
    diagram_fn=d_bothgusset,
    pro_tip="Unlike standup — gusset is on sides (width direction), not bottom. No gusset/2 in repeat. Higher film width per lane but simple repeat formula."
)

pouch_slide(prs, 8, "⑤", "3D Pouch / Flat Bottom Pouch", C_TEAL,
    "Premium block-bottom — stands upright with a flat rectangular base",
    film_f   = "Film Width  =  W  +  (2 × Side Gusset)",
    repeat_f = "Repeat  =  H  +  Top Seal  +  Bottom Gusset ÷ 2",
    example_calc = "W=180, H=280, Top=12, Side=35, BtmGusset=80 → Film=250mm, Repeat=332mm",
    uses=[
        "Premium coffee & specialty tea — high street retail",
        "Cosmetics & personal care — serums, powders",
        "Gourmet food — imported chocolate, artisan products",
        "Organic & natural — quinoa, chia, health foods",
        "Gift packaging — premium dry fruits, sweets",
    ],
    products=[
        "Ground coffee (250g)", "Specialty tea", "Premium chocolate",
        "Organic quinoa & seeds", "Turmeric latte powder", "Face wash powder",
        "Bath salts", "Protein isolate", "Premium dry fruits", "Herbal supplements",
    ],
    materials=[
        "Kraft paper + PE  (sustainable / eco look — coffee)",
        "PET + AL Foil + PE  (maximum barrier — coffee freshness)",
        "Matte OPP + CPP  (soft-touch premium look)",
        "PET + MET PET + PE  (metallic premium — cosmetics)",
    ],
    diagram_fn=d_flatbottom,
    pro_tip="Both side gusset AND bottom gusset — most complex to plan. Bottom gusset/2 in repeat (like standup). Side gussets in film width. Premium product only."
)

# ══════════════════════════════════════════════════════════════
# SLIDE 9  ── FILM STRUCTURE GUIDE
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl); top_bar(sl, C_GREEN)
section_tag(sl, "MATERIAL GUIDE", C_GREEN)
heading(sl, "Film Structure Selection Guide",
        "Which laminate to choose based on product type and barrier requirement", C_GREEN)

structs = [
    (C_BLUE,   "BOPP + CPP",
     "Low cost · Good clarity · Seals well",
     "Namkeen, biscuits, dry snacks, confectionery",
     "Low barrier — not for moisture-sensitive products"),
    (C_GREEN,  "PET + PE",
     "Good strength · Medium barrier · Food safe",
     "Milk, liquid pouches, oil (short shelf life)",
     "Standard choice for most flexible pouches"),
    (C_ORANGE, "PET + AL Foil + PE",
     "Maximum barrier · Long shelf life · Light-proof",
     "Coffee, spices, pharma, protein powder, chips",
     "Most expensive — use only when shelf life > 6 months"),
    (C_PURPLE, "PET + MET PET + PE",
     "High barrier · Metallic look · Premium feel",
     "Premium snacks, dry fruits, cosmetics",
     "Good balance of cost and barrier for premium segment"),
    (C_TEAL,   "Kraft Paper + PE",
     "Eco look · Biodegradable appearance · Sustainable branding",
     "Organic coffee, health foods, natural cosmetics",
     "Lower barrier than foil — needs dry product"),
    (C_YELLOW, "Woven PP Laminate",
     "Very high strength · Heavy loads · Tear resistant",
     "Rice (10kg+), fertilizer, cement, animal feed",
     "Industrial use only — not for food contact without liner"),
]
for i,(col,name,props,use,note) in enumerate(structs):
    cx = 0.5 + (i%2)*6.42
    cy = 1.8 + (i//2)*1.72
    rect(sl, cx, cy, 6.2, 1.6, C_CARD)
    rect(sl, cx, cy, 0.07, 1.6, col)
    txt(sl, name, cx+0.18, cy+0.08, 5.8, 0.36, sz=13, bold=True, col=col)
    txt(sl, props, cx+0.18, cy+0.44, 5.8, 0.28, sz=9, col=C_LIGHT, italic=True)
    txt(sl, "Use for: "+use, cx+0.18, cy+0.7, 5.8, 0.3, sz=9, col=C_WHITE)
    txt(sl, "⚠ "+note, cx+0.18, cy+0.98, 5.8, 0.3, sz=8, col=C_MUTED, italic=True)

slide_num(sl, 9); bot_bar(sl, C_GREEN)

# ══════════════════════════════════════════════════════════════
# SLIDE 10  ── CALCULATION EXAMPLE (3 SIDE SEAL LIVE WALKTHROUGH)
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl); top_bar(sl, C_ORANGE)
section_tag(sl, "LIVE EXAMPLE", C_ORANGE)
heading(sl, "Complete Calculation Walkthrough",
        "3 Side Seal Pouch — 130mm × 200mm — 1080mm machine", C_ORANGE)

steps = [
    ("STEP 1", "Enter Pouch Dimensions",   "Pouch Width (W) = 130 mm\nPouch Height (H) = 200 mm\nSide Seal = 10 mm  |  Top Seal = 10 mm  |  Bottom Seal = 8 mm", C_BLUE),
    ("STEP 2", "Calculate Film Width",     "Film Width = W + 2 × Side Seal\n= 130 + 2×10 = 150 mm per lane\n(this is what 1 AC UPS occupies across the film)", C_GREEN),
    ("STEP 3", "Calculate Repeat Length",  "Repeat = H + Top Seal + Bottom Seal\n= 200 + 10 + 8 = 218 mm\n(this must divide evenly into cylinder circumference)", C_ORANGE),
    ("STEP 4", "Calculate AC UPS",         "Machine Width = 1080 mm  |  Trim = 10mm each side\nUsable width = 1080 - 20 = 1060 mm\nAC UPS = floor(1060 ÷ 150) = 7 lanes", C_PURPLE),
    ("STEP 5", "Select Cylinder",          "Cylinder Circ must be multiple of 218 mm:\n218×1=218  |  218×2=436  |  218×3=654 mm\nChoose 436mm cylinder → 2 repeats per revolution", C_TEAL),
    ("STEP 6", "Total UPS & Film Usage",   "Total UPS = 7 AC UPS × 2 Repeat UPS = 14 pieces/rev\nFor 1,00,000 pouches: ceil(1,00,000 ÷ 14) = 7,143 meters\nFilm used: 7,143 m × 1060mm × GSM = X kg", C_YELLOW),
]
for i,(step,title,body,col) in enumerate(steps):
    cx = 0.5 + (i%2)*6.42
    cy = 1.82 + (i//2)*1.78
    rect(sl, cx, cy, 6.2, 1.68, C_CARD)
    rect(sl, cx, cy, 0.07, 1.68, col)
    rect(sl, cx+0.14, cy+0.1, 0.9, 0.34, col)
    txt(sl, step, cx+0.14, cy+0.1, 0.9, 0.34, sz=9, bold=True,
        col=C_BG, align=PP_ALIGN.CENTER)
    txt(sl, title, cx+1.14, cy+0.1, 4.9, 0.34, sz=12, bold=True, col=col)
    txt(sl, body,  cx+0.18, cy+0.52, 5.9, 1.1, sz=9, col=C_LIGHT, wrap=True)

slide_num(sl, 10); bot_bar(sl, C_ORANGE)

# ══════════════════════════════════════════════════════════════
# SLIDE 11  ── COMMON MISTAKES
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl); top_bar(sl, C_RED)
section_tag(sl, "AVOID THESE", C_RED)
heading(sl, "Common Mistakes in Pouch Planning",
        "Errors that cause machine stops, cylinder rejection, or film waste", C_RED)

mistakes = [
    ("❌ Forgetting to add seals to film width",
     "Seal width (8–12mm per side) MUST be added to W.\nForgetting adds 16–24mm error per lane → 1–2 UPS less than expected."),
    ("❌ Not adding gusset/2 to repeat length",
     "For standup & 3D pouches, gusset folds in half.\nMissing this = cylinder circumference is short → print misalign."),
    ("❌ Using pouch W as film width for center seal",
     "Center seal: W is HALF width. Film = 2W + seal.\nUsing just W gives 50% wrong film — machine cannot run."),
    ("❌ Confusing side gusset vs. bottom gusset",
     "Both-side gusset → gusset in width direction (adds to film W).\n3D Flat Bottom → gusset in both directions (adds to W and H)."),
    ("❌ Cylinder circ not a multiple of repeat",
     "Even 1mm off = print register error on every revolution.\nAlways verify: cyl_circ % repeat_length < 0.5mm."),
    ("❌ Wrong structure type for product",
     "Heavy product in 3-side seal = bag bursts on filling line.\nLiquid in standup without proper seal = leakage at gusset junction."),
]
for i,(title,desc) in enumerate(mistakes):
    cx = 0.5 + (i%2)*6.42
    cy = 1.78 + (i//2)*1.78
    rect(sl, cx, cy, 6.2, 1.68, C_CARD)
    rect(sl, cx, cy, 0.07, 1.68, C_RED)
    txt(sl, title, cx+0.18, cy+0.1, 5.9, 0.42, sz=11, bold=True,
        col=RGBColor(0xFC,0xA5,0xA5))
    txt(sl, desc, cx+0.18, cy+0.56, 5.9, 1.0, sz=10, col=C_LIGHT, wrap=True)

slide_num(sl, 11); bot_bar(sl, C_RED)

# ══════════════════════════════════════════════════════════════
# SLIDE 12  ── HOW AJ SHRINK ERP HANDLES THIS
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl); top_bar(sl, C_BLUE)
section_tag(sl, "OUR SOLUTION", C_BLUE)
heading(sl, "How AJ Shrink ERP Automates All This",
        "Select pouch type → enter dimensions → system calculates everything automatically")

features = [
    (C_BLUE,   "Auto Film Width",
     "Select pouch type and the system auto-applies the correct formula.\n"
     "3 Side = W+2S  |  Center = 2W+S  |  Gusset = W+2G\n"
     "No manual formula — zero calculation errors."),
    (C_GREEN,  "Auto Repeat Length",
     "Top seal, bottom seal, and gusset/2 are all added automatically.\n"
     "The system picks the correct formula per pouch type.\n"
     "Cylinder circumference is verified as exact multiple."),
    (C_ORANGE, "Live UPS Planning",
     "Enter machine width and trims — system shows all valid plans.\n"
     "AC UPS across × Repeat UPS = total pieces per revolution.\n"
     "Best plan highlighted automatically (least waste, most UPS)."),
    (C_PURPLE, "Cylinder Matching",
     "Real cylinder stock checked first (±1mm tolerance).\n"
     "If no stock match → special order cylinder auto-suggested.\n"
     "Sleeve-based or cylinder-direct planning — both supported."),
    (C_TEAL,   "Dimension Diagram",
     "Visual diagram shows exact film zones — seals, gussets, body.\n"
     "Instant feedback as you type dimensions.\n"
     "Client presentations — show exactly what will be printed."),
    (C_YELLOW, "Cost Estimation",
     "Film cost, ink cost, cylinder cost, machine time — all calculated.\n"
     "3 quantity breakdowns (Q1/Q2/Q3) in a single view.\n"
     "Export to catalog → order → work order in one flow."),
]
for i,(col,title,desc) in enumerate(features):
    cx = 0.5 + (i%2)*6.42
    cy = 1.78 + (i//2)*1.78
    rect(sl, cx, cy, 6.2, 1.68, C_CARD)
    rect(sl, cx, cy, 0.07, 1.68, col)
    rect(sl, cx+0.14, cy+0.12, 0.42, 0.42, col)
    txt(sl, str(i+1), cx+0.14, cy+0.12, 0.42, 0.42,
        sz=16, bold=True, col=C_BG, align=PP_ALIGN.CENTER)
    txt(sl, title,  cx+0.68, cy+0.12, 5.38, 0.4, sz=13, bold=True, col=col)
    txt(sl, desc,   cx+0.18, cy+0.58, 5.9,  1.0, sz=9,  col=C_LIGHT, wrap=True)

slide_num(sl, 12); bot_bar(sl, C_BLUE)

# ══════════════════════════════════════════════════════════════
# SLIDE 13  ── SUMMARY COMPARISON
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl); top_bar(sl, C_PURPLE)
section_tag(sl, "SUMMARY", C_PURPLE)
heading(sl, "Which Pouch for Which Product?",
        "Quick decision guide for sales & estimation teams", C_PURPLE)

data = [
    ("3 Side Seal",    C_BLUE,   "Low",     "Low–Medium","No",    "Flat only",   "Sachets, spices, seeds, snacks, pharma"),
    ("Center Seal",    C_GREEN,  "Low",     "Low",       "No",    "Flat / flow", "Milk, oil, water, sauce, candy, noodles"),
    ("Standup Zipper", C_ORANGE, "Medium",  "Medium",    "Yes",   "Stands up",   "Dry fruits, coffee, pet food, protein"),
    ("Both Gusset",    C_PURPLE, "High",    "Medium",    "No",    "Flat+gusset", "Rice, atta, fertilizer, bulk grains"),
    ("3D Flat Bottom", C_TEAL,   "Premium", "Medium",    "Optional","Block base","Premium coffee, tea, cosmetics, organic"),
]
hdrs = ["Pouch Type","Cost","Barrier Req","Resealable","How it Stands","Best For"]
col_x = [0.5, 2.65, 3.9,  5.3,  6.55, 7.8]
col_w = [2.05,1.15, 1.3,  1.15, 1.15, 5.3]
# header
rect(sl, 0.5, 1.78, 12.33, 0.44, RGBColor(0x10,0x1C,0x30))
for h,cx,cw in zip(hdrs, col_x, col_w):
    txt(sl, h, cx+0.06, 1.82, cw, 0.36, sz=9, bold=True,
        col=C_PURPLE, align=PP_ALIGN.LEFT)

for ri,(name,col,cost,barrier,reseal,stand,best) in enumerate(data):
    ry = 2.28 + ri*1.02
    fill = RGBColor(0x10,0x1C,0x30) if ri%2==0 else C_CARD
    rect(sl, 0.5, ry, 12.33, 0.96, fill)
    rect(sl, 0.5, ry, 0.07, 0.96, col)
    txt(sl, name,    col_x[0]+0.1, ry+0.08, col_w[0]-0.1, 0.8, sz=11, bold=True, col=col)
    txt(sl, cost,    col_x[1]+0.05,ry+0.28, col_w[1]-0.05, 0.4, sz=10, col=C_LIGHT, align=PP_ALIGN.CENTER)
    txt(sl, barrier, col_x[2]+0.05,ry+0.28, col_w[2]-0.05, 0.4, sz=10, col=C_LIGHT, align=PP_ALIGN.CENTER)
    txt(sl, reseal,  col_x[3]+0.05,ry+0.28, col_w[3]-0.05, 0.4, sz=10, col=C_LIGHT, align=PP_ALIGN.CENTER)
    txt(sl, stand,   col_x[4]+0.05,ry+0.28, col_w[4]-0.05, 0.4, sz=9,  col=C_MUTED, align=PP_ALIGN.CENTER)
    txt(sl, best,    col_x[5]+0.05,ry+0.08, col_w[5]-0.1,  0.8, sz=9,  col=C_LIGHT)

slide_num(sl, 13); bot_bar(sl, C_PURPLE)

# ══════════════════════════════════════════════════════════════
# SLIDE 14  ── THANK YOU
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
rect(sl, 0, 0, 13.33, 0.07, C_BLUE)
# glow rect center
rect(sl, 3.8, 1.6, 5.7, 4.4, C_CARD)
rect(sl, 3.8, 1.6, 5.7, 0.08, C_BLUE)
txt(sl, "AJ SHRINK ERP", 3.8, 1.82, 5.7, 0.55,
    sz=18, bold=True, col=C_BLUE, align=PP_ALIGN.CENTER)
txt(sl, "Thank You", 3.8, 2.45, 5.7, 0.85,
    sz=44, bold=True, col=C_WHITE, align=PP_ALIGN.CENTER)
rect(sl, 4.8, 3.4, 3.6, 0.06, C_BLUE)
txt(sl, "Flexible Packaging · Complete ERP Solution",
    3.8, 3.55, 5.7, 0.45, sz=13, col=C_MUTED, align=PP_ALIGN.CENTER)
txt(sl, "5 Pouch Types  ·  Auto Film Calculation  ·  Cylinder Planning  ·  Cost Estimation",
    3.8, 4.1,  5.7, 0.4, sz=10, col=C_BLUE, align=PP_ALIGN.CENTER, italic=True)
txt(sl, "Powered by AJ Shrink ERP",
    3.8, 4.65, 5.7, 0.35, sz=10, col=C_MUTED, align=PP_ALIGN.CENTER, italic=True)
# 5 dots
for i,(n,c) in enumerate(zip(["3SS","CS","SUP","BSG","3D"],POUCH_C)):
    xi = 4.1 + i*1.0
    rect(sl, xi, 5.6, 0.82, 0.45, c)
    txt(sl, n, xi, 5.6, 0.82, 0.45, sz=9, bold=True,
        col=C_BG, align=PP_ALIGN.CENTER)
bot_bar(sl, C_BLUE)

# ── Save ──────────────────────────────────────────────────────
out = "d:/AJ shrink new code/New updated code/Pouch_Types_Presentation_v2.pptx"
prs.save(out)
print(f"Saved: {out}")
